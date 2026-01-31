"""
PowerPoint (.pptx) content extractor.
"""

import re
from pathlib import Path
from typing import Optional

from .base import BaseExtractor, ExtractedContent


class PowerPointExtractor(BaseExtractor):
    """Extract content from PowerPoint presentations."""

    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def extract(self) -> ExtractedContent:
        """
        Extract content from PowerPoint file.

        Uses python-pptx for structured extraction.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint extraction. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation(str(self.file_path))

        slides = []
        all_text = []
        tables = []
        images = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = self._extract_slide_content(slide, slide_num)
            slides.append(slide_content)
            all_text.append(slide_content["content"])

            # Extract tables from this slide
            slide_tables = self._extract_tables(slide)
            tables.extend(slide_tables)

            # Track images
            slide_images = self._extract_images(slide, slide_num)
            images.extend(slide_images)

        # Extract metadata
        metadata = self._extract_metadata(prs)

        return ExtractedContent(
            raw_text="\n\n".join(all_text),
            slides=slides,
            tables=tables,
            images=images,
            metadata=metadata,
            source_file=str(self.file_path)
        )

    def _extract_slide_content(self, slide, slide_num: int) -> dict:
        """Extract all content from a single slide."""
        title = ""
        body_text = []
        notes = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = self._get_text_from_shape(shape)
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Check if it's a title placeholder
                    if ph_type in [1, 2, 3]:  # TITLE, CENTER_TITLE, SUBTITLE
                        title = text
                    else:
                        body_text.append(text)
                else:
                    body_text.append(text)

            # Handle tables
            if shape.has_table:
                table_text = self._table_to_text(shape.table)
                body_text.append(table_text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text

        content = title
        if body_text:
            content += "\n" + "\n".join(body_text)

        return {
            "number": slide_num,
            "title": title.strip(),
            "content": content.strip(),
            "body": "\n".join(body_text).strip(),
            "notes": notes.strip(),
            "has_table": any(shape.has_table for shape in slide.shapes if hasattr(shape, 'has_table')),
            "has_image": any(hasattr(shape, 'image') for shape in slide.shapes)
        }

    def _get_text_from_shape(self, shape) -> str:
        """Extract text from a shape with text frame."""
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs)
            if text.strip():
                # Preserve bullet structure
                level = paragraph.level
                prefix = "  " * level + "• " if level > 0 else ""
                paragraphs.append(prefix + text.strip())
        return "\n".join(paragraphs)

    def _table_to_text(self, table) -> str:
        """Convert a table to text format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_tables(self, slide) -> list[list[list[str]]]:
        """Extract all tables from a slide."""
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
        return tables

    def _extract_images(self, slide, slide_num: int) -> list[dict]:
        """Extract image metadata from a slide."""
        images = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                images.append({
                    "slide_number": slide_num,
                    "content_type": shape.image.content_type,
                    "size": len(shape.image.blob),
                })
        return images

    def _extract_metadata(self, prs) -> dict:
        """Extract presentation metadata."""
        core_props = prs.core_properties
        return {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "slide_count": len(prs.slides),
        }
